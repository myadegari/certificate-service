from pptx import Presentation
import comtypes.client
import os
import re
import jdatetime
import uuid
import qrcode
import asyncio
from concurrent.futures import ThreadPoolExecutor
from pptxtopdf import convert

def generate_qr_code(data: str, save_path: str):
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.ERROR_CORRECT_H,
        box_size=10,
        border=4,
    )
    qr.add_data(data)
    qr.make(fit=True)
    img = qr.make_image(fill_color="#DCDCDC", back_color="white")
    with open(save_path, "wb") as f:
        img.save(f)

def add_qr_to_placeholder(slide, qr_data, qr_output_dir="tmp"):
    if not os.path.exists(qr_output_dir):
        os.makedirs(qr_output_dir)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text
        image_placeholders = re.findall(r'{{qr:(.*?)}}', text)
        if image_placeholders:
            for img_key in image_placeholders:
                if img_key in qr_data:
                    data = qr_data[img_key]
                    qr_path = os.path.join(qr_output_dir, f"qr_{hash(data)}.png")
                    generate_qr_code(data, qr_path)
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    shape._element.getparent().remove(shape._element)
                    slide.shapes.add_picture(qr_path, left, top, width, height)
                    break

def replace_placeholder_preserve_style(shape, data):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for key, value in data.items():
                placeholder = f"{{{{{key}}}}}"
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, value)

def add_image_to_placeholder(slide, image_data):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text
        image_placeholders = re.findall(r'{{image:(.*?)}}', text)
        if image_placeholders:
            for img_key in image_placeholders:
                if img_key in image_data:
                    img_path = image_data[img_key]
                    if not os.path.exists(img_path):
                        img_path = "services/n-image.png"
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    shape._element.getparent().remove(shape._element)
                    slide.shapes.add_picture(img_path, left, top, width, height)
                    break

async def generate_certificate(template_path, output_dir, text_data, image_data=None, qr_data=None):
    if image_data is None:
        image_data = {}
    if qr_data is None:
        qr_data = {}
    prs = Presentation(template_path)
    slide = prs.slides[0]
    if image_data:
        add_image_to_placeholder(slide, image_data)
    if qr_data:
        add_qr_to_placeholder(slide, qr_data)
    for shape in slide.shapes:
        replace_placeholder_preserve_style(shape, text_data)
    output_filename = f"certificate_{text_data.get('unique', 'filled').replace(' ', '_')}"
    output_pptx = os.path.join(output_dir, f"{output_filename}.pptx")
    prs.save(output_pptx)
    return output_pptx

async def convert_to_pdf(pptx_path, output_dir):
    loop = asyncio.get_running_loop()
    def windows_convert():
        comtypes.CoInitialize()
        try:
            # powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            # Do not set Visible = False to avoid error
            try:
                print(f"Opening PPTX: {pptx_path}")
                # ppt = powerpoint.Presentations.Open(pptx_path, WithWindow=0)  # Minimize window
                base_filename = os.path.splitext(os.path.basename(pptx_path))[0]
                print(f"Generated PDF: {output_dir}")
                output_path = os.path.join(output_dir, f"{base_filename}.pdf")
                convert(pptx_path, output_dir)
                # ppt.SaveAs(output_path, 32)  # 32 = PDF
                # ppt.Close()
            except Exception as e:
                raise Exception(f"PowerPoint error: {str(e)}")
            # finally:
                # powerpoint.Quit()
        finally:
            comtypes.CoUninitialize()
        if os.path.exists(pptx_path):
            os.remove(pptx_path)  # Cleanup PPTX
        return output_path
    return await loop.run_in_executor(None, windows_convert)

def safe_image(image_path: str) -> str:
    return image_path if os.path.exists(image_path) else "services/n-image.png"